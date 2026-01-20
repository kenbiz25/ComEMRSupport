
# core/indexing/pipeline.py
import os
import pathlib
import logging
import time
import random
import hashlib
import unicodedata
import re
from typing import List, Iterable, Dict, Any, Optional, Sequence

import numpy as np
from openai import OpenAI
# NOTE: Depending on openai package version, these exceptions may reside under different modules.
# If your environment complains, you can catch a broad Exception in _with_backoff or import accordingly.
from openai import APIError, RateLimitError, InternalServerError  # type: ignore

from config.settings import settings
from core.knowledge.store_faiss import FaissStore

# Optional readers (install: python-docx, PyPDF2, python-pptx, tiktoken)
try:
    from docx import Document  # python-docx
except Exception:
    Document = None

try:
    from PyPDF2 import PdfReader  # PyPDF2
except Exception:
    PdfReader = None

try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None

try:
    import tiktoken  # for token-aware chunking
except Exception:
    tiktoken = None


# -----------------------------------------------------------------------------
# logging
# -----------------------------------------------------------------------------
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)


# -----------------------------------------------------------------------------
# helpers: model, normalization, chunking
# -----------------------------------------------------------------------------

def _is_openai_embed(model_name: str) -> bool:
    return str(model_name).startswith("text-embedding-")

def _embed_dim(model_name: str) -> int:
    if model_name == "text-embedding-3-small":
        return 1536
    if model_name == "text-embedding-3-large":
        return 3072
    # default
    return 1536

def _l2_normalize(vecs: Sequence[Sequence[float]]) -> List[List[float]]:
    out: List[List[float]] = []
    for v in vecs:
        arr = np.asarray(v, dtype=np.float32)
        norm = np.linalg.norm(arr)
        if norm == 0:
            out.append(arr.tolist())
        else:
            out.append((arr / norm).tolist())
    return out

_whitespace_re = re.compile(r"\s+")
def _normalize_whitespace(s: str) -> str:
    # collapse multiple whitespace into single space
    return _whitespace_re.sub(" ", s).strip()

def _clean_text(s: str) -> str:
    # Normalize unicode and whitespace to reduce duplicates/inconsistencies
    s = unicodedata.normalize("NFC", s or "")
    s = _normalize_whitespace(s)
    return s

def _chunk_by_chars(text: str, size: int = 1800, overlap: int = 200) -> List[str]:
    if not text:
        return []
    chunks: List[str] = []
    i, n = 0, len(text)
    step = max(size - overlap, 1)
    while i < n:
        chunks.append(text[i : i + size])
        i += step
    return chunks

def _chunk_by_tokens(text: str, model: str, chunk_size: int = 800, overlap: int = 100) -> List[str]:
    """
    Token-aware chunking using tiktoken; falls back to char-based if not available.
    """
    if not text:
        return []
    if tiktoken is None:
        # rough fallback: approximate 4 chars per token
        return _chunk_by_chars(text, size=chunk_size * 4, overlap=overlap * 4)

    try:
        # If model not recognized by tiktoken, use a close cousin
        try:
            enc = tiktoken.encoding_for_model(model)
        except Exception:
            enc = tiktoken.get_encoding("cl100k_base")
        toks = enc.encode(text)
        step = max(chunk_size - overlap, 1)
        chunks: List[str] = []
        for i in range(0, len(toks), step):
            piece = toks[i : i + chunk_size]
            if not piece:
                continue
            chunks.append(enc.decode(piece))
        return chunks
    except Exception:
        # defensive fallback
        return _chunk_by_chars(text, size=chunk_size * 4, overlap=overlap * 4)

# -----------------------------------------------------------------------------
# helpers: readers
# -----------------------------------------------------------------------------

def _read_docx(path: pathlib.Path) -> str:
    if Document is None:
        return ""
    try:
        doc = Document(str(path))
        parts: List[str] = []

        # Paragraphs
        parts.extend(p.text for p in doc.paragraphs if p.text)

        # Tables
        for table in getattr(doc, "tables", []):
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        parts.append(cell.text)

        return _clean_text("\n".join(parts))
    except Exception as e:
        logger.warning(f"DOCX parse failed for {path.name}: {e}")
        return ""

def _read_pdf(path: pathlib.Path) -> str:
    if PdfReader is None:
        return ""
    try:
        reader = PdfReader(str(path))
        buf: List[str] = []
        for page in getattr(reader, "pages", []):
            t = page.extract_text() or ""
            if t:
                buf.append(t)
        return _clean_text("\n".join(buf))
    except Exception as e:
        logger.warning(f"PDF parse failed for {path.name}: {e}")
        return ""

def _read_pptx(path: pathlib.Path) -> str:
    if Presentation is None:
        return ""
    try:
        prs = Presentation(str(path))
        buf: List[str] = []
        for slide in prs.slides:
            # Shapes with text
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text:
                    buf.append(shape.text)
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    for row in tbl.rows:
                        for cell in row.cells:
                            if cell.text:
                                buf.append(cell.text)
            # Notes
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text or ""
                if notes_text:
                    buf.append(notes_text)
        return _clean_text("\n".join(buf))
    except Exception as e:
        logger.warning(f"PPTX parse failed for {path.name}: {e}")
        return ""

def _read_text_like(path: pathlib.Path) -> str:
    try:
        return _clean_text(path.read_text(encoding="utf-8", errors="ignore"))
    except Exception as e:
        logger.warning(f"Text read failed for {path.name}: {e}")
        return ""


# -----------------------------------------------------------------------------
# helpers: OpenAI retries
# -----------------------------------------------------------------------------

def _with_backoff(callable_, *args, max_retries: int = 6, **kwargs):
    delay = 1.0
    for attempt in range(max_retries):
        try:
            return callable_(*args, **kwargs)
        except (RateLimitError, InternalServerError, APIError) as e:
            if attempt == max_retries - 1:
                raise
            sleep_for = delay + random.random() * 0.5
            logger.info(f"Retrying after error ({e}). attempt={attempt+1}/{max_retries}, sleep={sleep_for:.2f}s")
            time.sleep(sleep_for)
            delay = min(delay * 2, 30.0)


# -----------------------------------------------------------------------------
# pipeline
# -----------------------------------------------------------------------------

class IndexingPipeline:
    """
    Reindexes KB content using OpenAI embeddings and stores into FAISS.
    Streaming per file, with retries and cleanup.
    """

    SUPPORTED_EXTS = {".docx", ".pdf", ".pptx", ".txt", ".md"}

    def __init__(
        self,
        batch_size: int = 128,
        kb_dir: Optional[str] = None,
        namespace: Optional[str] = None,
        token_chunk_size: int = 800,
        token_overlap: int = 100,
        max_file_mb: int = 50,
    ):
        self.model_name = settings.EMBED_MODEL
        self.use_openai = _is_openai_embed(self.model_name)
        self.dim = _embed_dim(self.model_name)
        self.batch_size = batch_size
        self.token_chunk_size = token_chunk_size
        self.token_overlap = token_overlap
        self.max_file_bytes = max_file_mb * 1024 * 1024

        self.kb_dir = pathlib.Path(kb_dir or settings.KB_DIR)
        self.namespace = namespace or settings.KB_NAMESPACE

        # IMPORTANT: FaissStore should use IndexFlatIP / METRIC_INNER_PRODUCT if using L2-normalized vectors
        self.store = FaissStore(dim=self.dim, namespace=self.namespace)

        if self.use_openai:
            if not settings.OPENAI_API_KEY:
                raise RuntimeError("OPENAI_API_KEY is not set")
            # optional: pass organization if your org requires it
            client_kwargs = {"api_key": settings.OPENAI_API_KEY}
            if getattr(settings, "OPENAI_ORG", None):
                client_kwargs["organization"] = settings.OPENAI_ORG
            # add timeout to avoid indefinite hangs
            client_kwargs["timeout"] = getattr(settings, "OPENAI_TIMEOUT", 60)
            self.oa = OpenAI(**client_kwargs)
        else:
            raise RuntimeError(f"Unsupported embedding model: {self.model_name}")

    # -------------------------------------------------------------------------
    # chunking / parsing
    # -------------------------------------------------------------------------

    def _extract_text(self, path: pathlib.Path) -> str:
        ext = path.suffix.lower()
        if ext == ".docx":
            return _read_docx(path)
        if ext == ".pdf":
            return _read_pdf(path)
        if ext == ".pptx":
            return _read_pptx(path)
        if ext in {".txt", ".md"}:
            return _read_text_like(path)
        return ""

    def _iter_kb_files(self, root: pathlib.Path) -> Iterable[pathlib.Path]:
        if not root.exists():
            logger.warning(f"KB dir not found: {root}")
            return
        for p in root.glob("**/*"):
            if not p.is_file():
                continue
            # skip hidden dirs/files and temp files
            if any(seg.startswith(".") for seg in p.parts):
                continue
            if p.name.startswith(("~$", "._")):
                continue
            if p.suffix.lower() not in self.SUPPORTED_EXTS:
                continue
            try:
                if p.stat().st_size > self.max_file_bytes:
                    logger.warning(f"Skipping large file (> {self.max_file_bytes} bytes): {p}")
                    continue
            except Exception:
                # if stat fails, skip
                continue
            yield p

    # -------------------------------------------------------------------------
    # embeddings
    # -------------------------------------------------------------------------

    def _embed_batch_openai(self, texts: List[str]) -> List[List[float]]:
        resp = _with_backoff(
            self.oa.embeddings.create,
            model=self.model_name,
            input=texts,
        )
        vecs = [d.embedding for d in resp.data]
        # normalize to use cosine similarity with FAISS IP index
        return _l2_normalize(vecs)

    def _embed_in_batches(self, texts: List[str], batch_size: int) -> List[List[float]]:
        all_vecs: List[List[float]] = []
        for i in range(0, len(texts), batch_size):
            batch = texts[i : i + batch_size]
            embeddings = self._embed_batch_openai(batch)
            all_vecs.extend(embeddings)
        return all_vecs

    # -------------------------------------------------------------------------
    # storage helpers
    # -------------------------------------------------------------------------

    def _document_id_for(self, path: pathlib.Path) -> str:
        try:
            st = path.stat()
            key = f"{path.resolve()}::{int(st.st_mtime)}::{st.st_size}"
        except Exception:
            # fallback: path only
            key = f"{path.resolve()}::fallback"
        return hashlib.sha1(key.encode("utf-8")).hexdigest()

    def _upsert_batch(self, metas: List[Dict[str, Any]], embeddings: List[List[float]]) -> int:
        count = 0
        for m, vec in zip(metas, embeddings):
            meta = {
                "title": m["title"],
                "source_path": m["source_path"],
                "chunk_id": m["chunk_id"],
                "namespace": self.namespace,
                "document_id": m["document_id"],
            }
            self.store.upsert(m["text"], meta, vec)
            count += 1
        return count

    # -------------------------------------------------------------------------
    # per-file processing
    # -------------------------------------------------------------------------

    def _process_file(self, path: pathlib.Path) -> int:
        text = self._extract_text(path)
        if not text:
            logger.warning(f"Empty/failed parse: {path.name}")
            return 0

        # Token-aware chunking with fallback
        pieces = _chunk_by_tokens(text, self.model_name, self.token_chunk_size, self.token_overlap)
        if not pieces:
            return 0

        document_id = self._document_id_for(path)

        # Clean previous records for this source path or document id (depends on your FaissStore API)
        # Prefer a stable key like source_path; or document_id if your store supports versioning.
        try:
            if hasattr(self.store, "delete_by_source_path"):
                self.store.delete_by_source_path(self.namespace, str(path))
            elif hasattr(self.store, "delete_by_document_id"):
                self.store.delete_by_document_id(self.namespace, document_id)
        except Exception as e:
            logger.info(f"Cleanup skipped/failed for {path.name}: {e}")

        total = 0
        # Stream in batches
        for i in range(0, len(pieces), self.batch_size):
            batch = pieces[i : i + self.batch_size]
            metas = [
                {
                    "text": b,
                    "title": path.stem,
                    "source_path": str(path),
                    "chunk_id": i + j + 1,  # 1-based within file
                    "document_id": document_id,
                }
                for j, b in enumerate(batch)
            ]
            embeddings = self._embed_batch_openai([m["text"] for m in metas])
            total += self._upsert_batch(metas, embeddings)

        logger.info(f"File processed: {path.name} | chunks={total}")
        return total

    # -------------------------------------------------------------------------
    # public API
    # -------------------------------------------------------------------------

    def reindex_all(self) -> int:
        root = self.kb_dir
        files = list(self._iter_kb_files(root))
        logger.info(f"KB_DIR={root} | namespace={self.namespace} | files={len(files)}")
        if not files:
            return 0

        total_chunks = 0
        for path in files:
            try:
                total_chunks += self._process_file(path)
            except Exception as e:
                logger.exception(f"Failed processing {path}: {e}")

        # Save FAISS index to disk (atomic if supported)
        
        try:
            if hasattr(self.store, "save_atomic"):
                self.store.save_atomic()
            else:
                self.store.save()

            logger.info(
                f"FAISS index saved to {getattr(self.store, 'index_dir', '<unknown>')}"
            )

        except Exception as e:
            logger.error(f"Failed to save FAISS index: {e}")

